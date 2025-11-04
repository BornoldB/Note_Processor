import os
from pptx import Presentation

def extract_text_from_pptx(directory_path, output_directory):
    """
    Extract text from all PowerPoint (.pptx) files in a directory and save as text files.
    
    Args:
        directory_path (str): Directory containing PowerPoint files
        output_directory (str): Directory to save extracted text files
    """
    # Create output directory if it doesn't exist
    if not os.path.exists(output_directory):
        os.makedirs(output_directory)
    
    # Check if input directory exists
    if not os.path.exists(directory_path):
        print(f"Input directory '{directory_path}' does not exist.")
        return
    
    # Iterate through all files in the directory
    pptx_files = [f for f in os.listdir(directory_path) if f.lower().endswith('.pptx')]
    
    if not pptx_files:
        print(f"No .pptx files found in '{directory_path}' directory.")
        return
    
    for filename in pptx_files:
        pptx_path = os.path.join(directory_path, filename)
        output_filename = os.path.splitext(filename)[0] + '.txt'
        output_path = os.path.join(output_directory, output_filename)
        
        try:
            # Open and extract text from PowerPoint
            presentation = Presentation(pptx_path)
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = f"=== SLIDE {slide_num} ===\n"
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                # Only add slide if it has content
                if slide_text.strip() != f"=== SLIDE {slide_num} ===":
                    text_content.append(slide_text)
                else:
                    text_content.append(f"=== SLIDE {slide_num} ===\n[No text content]\n")
            
            # Combine all slide text
            full_text = "\n".join(text_content)
            
            # Save extracted text to a .txt file
            with open(output_path, 'w', encoding='utf-8') as txt_file:
                txt_file.write(full_text)
            
            print(f"Successfully extracted text from {filename} to {output_filename}")
            print(f"  - Processed {len(presentation.slides)} slides")
            
        except Exception as e:
            print(f"Error processing {filename}: {str(e)}")


def extract_text_from_single_pptx(pptx_path, output_path=None):
    """
    Extract text from a single PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_path (str): Optional output path for the text file
    
    Returns:
        str: Extracted text content
    """
    try:
        presentation = Presentation(pptx_path)
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = f"=== SLIDE {slide_num} ===\n"
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            # Only add slide if it has content
            if slide_text.strip() != f"=== SLIDE {slide_num} ===":
                text_content.append(slide_text)
            else:
                text_content.append(f"=== SLIDE {slide_num} ===\n[No text content]\n")
        
        # Combine all slide text
        full_text = "\n".join(text_content)
        
        # Save to file if output path is provided
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as txt_file:
                txt_file.write(full_text)
            print(f"Text extracted and saved to {output_path}")
        
        return full_text
        
    except Exception as e:
        print(f"Error processing {pptx_path}: {str(e)}")
        return ""


if __name__ == "__main__":
    # Specify the input and output directories
    input_directory = "pptxs"  # Directory containing PowerPoint files
    output_directory = "text_output"  # Directory to save text files
    
    print("Starting PowerPoint text extraction...")
    extract_text_from_pptx(input_directory, output_directory)
    print("Text extraction completed!")